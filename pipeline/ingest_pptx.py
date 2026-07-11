"""
Tier 1 (mechanical extraction, no AI judgment).

Pulls title, bullet text, and speaker notes out of a .pptx file using
python-pptx. Speaker notes often carry the presenter's fuller explanation,
so they're kept as their own field rather than merged into the bullets.
"""

import json
import sys

from pptx import Presentation


def extract_slides(pptx_path: str) -> list[dict]:
    prs = Presentation(pptx_path)
    slides = []

    for i, slide in enumerate(prs.slides):
        title = ""
        bullets = []

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            text = shape.text_frame.text.strip()
            if not text:
                continue

            is_title = shape == slide.shapes.title
            if is_title and not title:
                title = text
                continue

            for line in text.split("\n"):
                line = line.strip()
                if line:
                    bullets.append(line)

        notes = ""
        if slide.has_notes_slide:
            notes = (slide.notes_slide.notes_text_frame.text or "").strip()

        if not title and not bullets:
            continue  # skip blank/section-divider slides

        slides.append(
            {
                "slide_index": i,
                "title": title or f"Slide {i + 1}",
                "bullets": bullets,
                "notes": notes,
            }
        )

    return slides


def main():
    if len(sys.argv) < 2:
        print("usage: python ingest_pptx.py <deck.pptx> [output.json]", file=sys.stderr)
        sys.exit(1)

    pptx_path = sys.argv[1]
    slides = extract_slides(pptx_path)

    output = json.dumps(slides, indent=2)
    if len(sys.argv) >= 3:
        with open(sys.argv[2], "w") as f:
            f.write(output)
        print(f"Extracted {len(slides)} slides -> {sys.argv[2]}")
    else:
        print(output)


if __name__ == "__main__":
    main()
