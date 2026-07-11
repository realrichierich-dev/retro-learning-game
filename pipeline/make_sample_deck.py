"""
Generates a small sample .pptx so the pipeline can be proven end-to-end
without Rich needing to supply a real deck first. Topic: basic astronomy,
picked as a neutral placeholder -- swap in a real deck any time.
"""

import sys

from pptx import Presentation
from pptx.util import Inches

SLIDES = [
    {
        "title": "The Solar System",
        "bullets": [
            "The Solar System has eight planets orbiting the Sun.",
            "The Sun contains about 99.8% of the Solar System's total mass.",
        ],
        "notes": "The Sun's gravity is what holds the entire Solar System together.",
    },
    {
        "title": "Mercury",
        "bullets": [
            "Mercury is the smallest planet and the closest to the Sun.",
            "A year on Mercury lasts only 88 Earth days.",
        ],
        "notes": "Mercury has almost no atmosphere, so its surface temperature swings hundreds of degrees between day and night.",
    },
    {
        "title": "Venus",
        "bullets": [
            "Venus is the hottest planet in the Solar System.",
            "Venus rotates in the opposite direction to most other planets.",
        ],
        "notes": "Venus's thick carbon dioxide atmosphere traps heat through a runaway greenhouse effect.",
    },
    {
        "title": "Earth",
        "bullets": [
            "Earth is the only known planet with liquid water on its surface.",
            "Earth's atmosphere is about 78% nitrogen and 21% oxygen.",
        ],
        "notes": "Earth's single large moon helps stabilize its axial tilt, which keeps seasons relatively steady.",
    },
    {
        "title": "Mars",
        "bullets": [
            "Mars is known as the Red Planet because of iron oxide on its surface.",
            "Mars has the largest volcano in the Solar System, Olympus Mons.",
        ],
        "notes": "Olympus Mons is about three times the height of Mount Everest.",
    },
    {
        "title": "Jupiter",
        "bullets": [
            "Jupiter is the largest planet in the Solar System.",
            "Jupiter's Great Red Spot is a giant storm larger than Earth.",
        ],
        "notes": "The Great Red Spot has been observed for at least 350 years.",
    },
]


def build(output_path: str):
    prs = Presentation()
    layout = prs.slide_layouts[1]  # title + content

    for slide_data in SLIDES:
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = slide_data["title"]

        body = slide.placeholders[1].text_frame
        body.clear()
        for i, bullet in enumerate(slide_data["bullets"]):
            p = body.paragraphs[0] if i == 0 else body.add_paragraph()
            p.text = bullet

        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = slide_data["notes"]

    prs.save(output_path)
    print(f"Wrote sample deck -> {output_path}")


if __name__ == "__main__":
    output_path = sys.argv[1] if len(sys.argv) > 1 else "sample_deck.pptx"
    build(output_path)
